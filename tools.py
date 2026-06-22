import re
import json
import time
import base64
from uuid import uuid4
import requests

from typing import Dict, List, Any
import uuid


from click import prompt
from groq import Groq
from config import client_groq, local_client, gemini_client
from schema import State
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
import os
import json
from playwright.sync_api import sync_playwright
import json
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# summary_file = Path(__file__).parent / "cache" / "summary.json"


# with open(summary_file, "r", encoding="utf-8") as f:
#         summary_data = json.load(f)
DEV_MODE = True



class Parsing_data:

    def __init__(self, raw_data):

        self.raw_content = raw_data

        self.image_dir = "saved_images"

        os.makedirs(self.image_dir, exist_ok=True)

        self.image_registry = {}

        self.final_list = self.formate_data()

    def download_web_image(self, url, output_dir="output_images"):
        """
        Downloads an image from a web URL and saves it to a local directory.
        Includes explicit browser headers to bypass server 403 blocks.
        """
        os.makedirs(output_dir, exist_ok=True)
        
        # Define a clean, modern browser User-Agent header
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36"
        }
        
        try:
            # Pass the headers into your request block
            response = requests.get(url, headers=headers, timeout=10, stream=True)
            
            if response.status_code != 200:
                print(f"Failed to download image from URL: {url}. Status code: {response.status_code}")
                return None
                
            path_without_query = url.split("?")[0]
            extension = os.path.splitext(path_without_query)[1].lower()
            
            if extension not in [".jpg", ".jpeg", ".png", ".webp", ".gif"]:
                content_type = response.headers.get("Content-Type", "")
                if "image/jpeg" in content_type:
                    extension = ".jpg"
                elif "image/png" in content_type:
                    extension = ".png"
                elif "image/webp" in content_type:
                    extension = ".webp"
                elif "image/gif" in content_type:
                    extension = ".gif"
                else:
                    extension = ".png"
                    
            image_id = f"img_{uuid.uuid4().hex[:8]}"
            filename = f"{image_id}{extension}"
            file_path = os.path.join(output_dir, filename)
            
            with open(file_path, "wb") as f:
                for chunk in response.iter_content(chunk_size=8192):
                    if chunk:
                        f.write(chunk)
                        
            return {
                "image_id": image_id,
                "image_path": file_path,
                "caption": "Downloaded Web Asset"
            }
        
        except Exception as e:
            print(f"Exception encountered during image download from {url}: {str(e)}")
            return None

    def save_base64_image(self, image_string):

        try:

            # Remove metadata
            header, encoded = image_string.split(",", 1)

            # Detect extension
            if "png" in header:
                ext = "png"
            elif "jpeg" in header or "jpg" in header:
                ext = "jpg"
            else:
                ext = "png"

            image_id = f"img_{uuid4().hex[:8]}"

            image_path = os.path.join(
                self.image_dir,
                f"{image_id}.{ext}"
            )

            image_data = base64.b64decode(encoded)

            with open(image_path, "wb") as f:
                f.write(image_data)

            return {
                "image_id": image_id,
                "image_path": image_path,
                "caption": "Extracted document image"
            }

        except Exception as e:

            print(f"Image decode error: {e}")

            return None

    # =====================================================
    # FORMAT DATA
    # =====================================================
    def formate_data(self):

        pattern = r"^(#{1,6})\s+(.+)$"

        matches = list(
            re.finditer(
                pattern,
                self.raw_content,
                flags=re.MULTILINE
            )
        )

        sections = []

        for i, match in enumerate(matches):

            heading_marks = match.group(1)

            heading_text = match.group(2).strip()

            start = match.end()

            end = (
                matches[i + 1].start()
                if i + 1 < len(matches)
                else len(self.raw_content)
            )

            content = self.raw_content[start:end].strip()

            # =====================================================
            # TABLE EXTRACTION
            # =====================================================
            table_pattern = r"(\|.+\|\n\|[-:\s|]+\|\n(?:\|.*\|\n?)*)"

            tables = re.findall(
                table_pattern,
                content,
                flags=re.MULTILINE
            )

            # import pdb;pdb.set_trace()
            # =====================================================
            # IMAGE EXTRACTION
            # =====================================================
            image_pattern = r"!\[.*?\]\((.*?)\)"

            raw_images = re.findall(
                image_pattern,
                content
            )

            processed_images = []

            for image_string in raw_images:

                # Only process base64 images
                if image_string.startswith("data:image"):

                    image_info = self.save_base64_image(
                        image_string
                    )

                    if image_info:
                        processed_images.append(image_info)
                        self.image_registry[image_info["image_id"]] = image_info

                elif image_string.startswith(("http://", "https://")):
                    # Check if the URL has already been processed to save network bandwidth
                    existing_match = next(
                        (img for img in processed_images if img.get("image_url_source") == image_string), 
                        None
                    )
                    if existing_match:
                        # Cache Hit: Reference the existing dict pointer. 
                        # Crucial: DO NOT append it to processed_images again.
                        image_info = existing_match
                    else:
                        # Cache Miss: Trigger a fresh stream download to disk
                        image_info = self.download_web_image(image_string)
                        if image_info:
                            # Stamp the source URL on the dictionary object for tracking
                            image_info["image_url_source"] = image_string
                            
                            # Append to state tracking collections exactly once
                            processed_images.append(image_info)
                            self.image_registry[image_info["image_id"]] = image_info

            # =====================================================
            # SKIP TOC
            # =====================================================
            if heading_text.lower() == "table of contents":
                continue

            sections.append({

                "id": len(sections) + 1,

                "heading": heading_text,

                "level": len(heading_marks),

                "content": content,

                "word_count": len(
                    content.split()
                ),

                "tables": tables,

                "images": processed_images
            })

        return sections

def parsing_raw_data(state: State):


    # if DEV_MODE and os.path.exists("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/parse.json"): 
    #      with open("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/parse.json") as f:
    #           return {"parsed_content": json.load(f)}
         
    parser = Parsing_data(state["raw_file_content"])

    # with open("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/parse.json", "w") as f:
    #     json.dump(parser.final_list, f)

    print(f"parsing result : {json.dumps(parser.final_list)}")

    return {
        "parsed_content": parser.final_list
    }


def designer_node(state: State):
    # source_context = [s["heading"] for s in state.get("summarized_content", [])]
    # We only send the main headings to the small LLM to save tokens
    context = [s["slide_title"] for s in state["slide_plan"]]

    # if DEV_MODE and os.path.exists("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/design.json"): 
    #         try: 
    #             with open("./home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/design.json") as f:
    #                 return {"design_content": json.load(f)}
    #         except Exception as e:
    #             print("Invalid cache file. Regenerating...")
            

    
    prompt = f"""
    You are an expert Senior Brand Designer and Creative Director specialized in corporate typography and visual identity.

Your task is to analyze the slide headings for a presentation and engineer a unique, high-impact visual design system tailored specifically to the subject matter's industry, tone, and emotional context.

INPUT HEADINGS:
{json.dumps(context)}

DESIGN ASSIGNMENT CONSTRAINTS:

1. ARCHETYPE SELECTION:
   Analyze the vocabulary of the headings. Assign exactly ONE design archetype from this matrix:
   - "Technical_Modern": For data, engineering, forensics, AI, or product development. (Prefers clean neo-grotesque or monospace styling).
   - "Executive_Corporate": For financial reports, board meetings, legal summaries, or operations. (Prefers high-end serif headers and clean sans-serif body text).
   - "Creative_Vibrant": For marketing pitches, design reviews, storytelling, or human-centric topics. (Prefers bold, expressive display typefaces).

2. BACKGROUND STYLE & VARIETY:
   To ensure slides do not look flat and identical, select a background texture:
   - "solid": Standard flat color execution.
   - "subtle_gradient": A clean 135-degree gradient blending the canvas background into a slightly deeper shade.
   - "tech_grid": A clean, ultra-low opacity engineering cross-grid lines pattern overlay (highly recommended for Technical_Modern).

3. COLOR HARMONY (THE 60-30-10 RULE):
   - Choose a base "theme" ('light' or 'dark') based on what fits the topic best.
   - "bg_canvas": The dominant color (60%). If theme is light. If dark, use deep slate/obsidian. Never use pure #000000 or pure #FFFFFF.
   - "primary": Main typography color for titles. Must have an absolute minimum contrast ratio of 4.5:1 against the bg_canvas.
   - "secondary": Color for subheaders, borders, and secondary text blocks (30%).
   - "accent": A bright, vibrant, or high-impact color reserved strictly for key visual anchors, active metrics, or layout icons (10%). Avoid bleeding neon colors.

4. DYNAMIC VARIATION TOKEN:
   - "inverted_section_bg": Provide a contrasting color to invert the theme completely for Section Dividers or Conclusion slides to break up visual fatigue.

5. TEXT COLOR:
    - "CHOOSE THE TEXT COLOR ACCORDING TO THE BACKGROUND COLOR" 
    - Do not use dark color on the dark background as it will make the text hard to read
------------------EVERY DESIGN SHOULD BE BASED ON THE TOPIC OF THE PPT-----------------------------------

OUTPUT FORMAT:
You must return ONLY a valid JSON object. Do not include markdown code wrappers, prose, introduction, or post-processing explanations.

REQUIRED JSON SCHEMA: {{
  "archetype": "Technical_Modern | Executive_Corporate | Creative_Vibrant",
  "theme": "light | dark",
  "background_style": "solid | subtle_gradient | tech_grid",
  "bg_canvas": "#HEXCODE",
  "primary": "#HEXCODE",
  "secondary": "#HEXCODE",
  "accent": "#HEXCODE",
  "inverted_section_bg": "#HEXCODE"
}}
    """
    
    # Call your small LLM here
    chat_completion = client_groq.chat.completions.create(
            messages = [
                    {
                            "role" : "user",
                            "content" : prompt
                }
            ],
            model="llama-3.3-70b-versatile",
    )
    response = chat_completion.choices[0].message.content
    try:

        clean_json = (
            response
            .replace("```json", "")
            .replace("```", "")
            .strip()
        )

        design_config = json.loads(clean_json)
        print(f"{design_config}")

    except Exception as e:

        print(f"Designer Node Error: {e}")

        # fallback theme
        design_config = {
            "primary": "#1F3A5F",
            "secondary": "#4F81BD",
            "accent": "#E67E22",
            "theme": "light"
        }

    

    # with open("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/design.json", "w") as f:
    #     json.dump(
    #         design_config,
    #         f,
    #         indent=2,
    #         ensure_ascii=False
    #     )
    return {
        "design_config": design_config
    }


def summary_helper(section: dict):
        
        prompt = f"""
                        You are an executive presentation writer.

                        Convert the following content into concise boardroom-style bullet points.

                        RULES:
                        - Maximum 5 bullets
                        - Maximum 15 words per bullet
                        - Preserve factual meaning
                        - No explanations
                        - No markdown headers
                        - Do not repeat the same content again and again
                        - Start every bullet with '-'
                        - Except pointer donot add anything else i the summary content.

                        CONTENT:
                        {section['content']}
                """
        response = local_client.chat(
                model='gemma3:latest', # or whichever model you downloaded
                messages=[{'role': 'user', 'content': prompt}]
        )
        
        raw_summary = response['message']['content'].strip()

        summary = []

        for line in raw_summary.split("\n"):

            line = line.strip()

            if not line:
                continue

            if line.startswith("-"):
                line = line[1:].strip()

            summary.append(line)
        print("=====================================================")
        print(summary)

        return summary



def summary_node(state: State):
        
        # if DEV_MODE and os.path.exists("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/summary.json"):
        #     try: 
        #         with open("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/summary.json") as f:
        #             return {"summarized_content": json.load(f)}
        #     except Exception as e:
        #         print("Invalid cache file. Regenerating...")
         
        print("*******************************************")
        print(f"parser: {json.dumps(state["parsed_content"], indent= 2)}")
        sections = state["parsed_content"]
        print(f"section: {sections}")
        summarized_sections = []

        for section in sections:

                if section["word_count"] > 50:
                        summary = summary_helper(section)
                else:
                        lines = re.split(
                        r"\n|- ",
                        section["content"]
                        )

                        summary = [
                            l.strip()
                            for l in lines
                            if l.strip()
                        ]

                summarized_sections.append({
                "heading": section["heading"],
                "summary": summary,
                "tables": section["tables"],
                "images": section["images"]
                })
        
        # with open("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/summary.json", "w") as f:
        #     json.dump(
        #         summarized_sections,
        #         f,
        #         indent=2,
        #         ensure_ascii=False
        #     )
        

        return {
                "summarized_content": summarized_sections
        }
# def planner_prompt(section):
        
#         prompt =f"""
#         You are a Senior Strategy Consultant. Your task is to take the provided acquisition summaries and organize them into a 10-15 slide executive presentation for a Board of Directors
#         Input Data: {new_data}
#         Your Requirements:
#                 1. Narrative Flow: Organize the slides logically (e.g., Strategy Overview -> Tech Domain Deep-dives -> Financial ROI -> Future Outlook).
#                 2. Executive Titles: Create punchy, active-voice titles (e.g., 'Scaling AI Leadership' instead of 'AI Acquisitions').
#                 3. Layout Selection: Assign one of these three layouts to each slide based on content density: 'Title_Only', 'Bullet_Slide', or 'Two_Column'.
#                 4. Constraint: Use the specific figures found in the data, such as the $5.9B AI bookings or the $865M restructuring cost.
                
#         Output Format: You MUST return ONLY a valid JSON array of objects. Do not include prose or explanations.
#         Schema Example:
#                 [
#                 {{
#                 "slide_number": 1,
#                 "title": "Slide Title Here",
#                 "content": ["Bullet 1", "Bullet 2"],
#                 "layout": "Bullet_Slide"
#                 }}
#                 ]"
#         """
#         response = gemini_client.models.generate_content(
#                 model="gemini-2.5-flash-lite", 
#                 contents=prompt
#         )

#         print("++++++++++++++++++++++++++++++++++++++++++++++++++")
#         print(response.text)

def get_section_by_heading(sections, source_heading):

    for section in sections:

        if (
            section.get("heading", "").strip().lower()
            ==
            source_heading.strip().lower()
        ):

            return section

    return {}

def repair_json_output(raw_text):

    # Remove markdown wrappers
    raw_text = raw_text.replace("```json", "").replace("```", "").strip()

    # Remove smart quotes
    raw_text = raw_text.replace("“", '"').replace("”", '"')
    raw_text = raw_text.replace("‘", "'").replace("’", "'")

    # Remove control chars
    raw_text = re.sub(r'[\x00-\x1F\x7F]', '', raw_text)

    # Escape invalid backslashes
    raw_text = re.sub(r'\\(?!["\\/bfnrtu])', r'\\\\', raw_text)

    # Fix accidental unescaped quotes inside values
    raw_text = re.sub(r'(?<!\\)"(?=[^,:{}\[\]]*")layout_plan', r'\"', raw_text)

    return raw_text

def planner_node(state: State):

    sections = state["summarized_content"]

    safe_sections = json.dumps(
        sections,
        ensure_ascii=False
    )

    prompt = f"""
You are a senior executive presentation designer.

You will receive structured business content in JSON format.

INPUT:
{safe_sections}

LAYOUT OPTIONS:
- Title_Only
- Full_Content
- Two_Column
- Visual_Focus

IMPORTANT:
Layout defines visual arrangement only.
A slide may contain any combination of:
- content
- table
- image
Layout's is decided on the basis of the content of the in the slide plan.
- if slide plan only contain title no content then it should be title only slide.

TASK:
Design a professional slide plan.

RULES:
- Maximum 15 slides
- Use fewer slides if possible
- Preserve source content exactly
- Do NOT rewrite content
- Do NOT modify tables
- Do NOT modify images
- Return ONLY valid JSON array
- No markdown
- No explanations

CRITICAL:

If source section contains image:
copy image field exactly.

If source section contains table:
copy table field exactly.

Never remove image.
Never remove table.

If image exists:
layout must be
Image_Only
or
Two_Column

If table exists:
layout must be
Table_Only
or
Two_Column

OUTPUT FORMAT:
[
  {{
    "slide_number": 1,
    "title": "string",
    "source_section": "string",
    "layout": "Title_Only | Bullet_Slide | Two_Column | Table_Only | Image_Only",
    "content": [],
    "table": [],
    "image": []
  }}
]
"""

    response = gemini_client.models.generate_content(
        model="gemini-3.1-flash-lite",
        contents=prompt
    )

    raw_text = response.text.strip()

    safe_json = repair_json_output(raw_text)

    try:
        planned_slides = json.loads(safe_json)

    except Exception as e:

        print("PLANNER JSON ERROR:", e)
        print("RAW RESPONSE:")
        print(raw_text)

        # ---------------- FALLBACK ----------------
        planned_slides = []

        for idx, section in enumerate(sections[:15], start=1):

            planned_slides.append({
                "slide_number": idx,
                "title": section.get(
                    "heading",
                    f"Slide {idx}"
                ),

                "source_section": section.get(
                    "heading",
                    f"Slide {idx}"
                ),

                "layout": (
                    "Table_Only"
                    if section.get("tables")
                    else "Bullet_Slide"
                ),

                "content": (
                    section.get("summary", [])
                    if isinstance(
                        section.get("summary"),
                        list
                    )
                    else [section.get("summary", "")]
                ),

                "table": section.get(
                    "tables",
                    []
                ),

                "image": section.get(
                    "images",
                    []
                )
            })

    return {
        "slide_plan": planned_slides
    }




def story_planner(summary):
    print(f"SUMMARY:{json.dumps(summary, indent=2)}")
    sections = summary

    safe_sections = json.dumps(
        sections,
        ensure_ascii=False
    )

    prompt = f"""
        You are a senior management consultant and executive presentation strategist.

Your task is to build a compelling presentation narrative from the provided content.

INPUT:
{safe_sections}

OBJECTIVE:
Create a logical executive storyline suitable for a Board of Directors presentation.

GUIDELINES:
1. Think like a strategy consultant.
2. Organize information into a clear narrative flow.
3. Add the content, table and the image a it in the slide
4. Group related sections together and preserve all important facts.
5. Maximum 15 slides. Fewer slides are preferred if the story remains complete.
6. Always remember to add a Conclusion Slide at the end.
7. Remember you don't have to summarize the content. You just need to set the story of the ppt and divide the data into 15 slides.

CRITICAL REQUIREMENT:
The final slide of the array MUST be a strategic conclusion, wrap-up, or forward-looking takeaway slide that sums up the core impact or next steps.

RULES:
- Do NOT create layouts or design instructions.
- Preserve source text context accurately inside the content array.
- Return ONLY a valid JSON array. No markdown blocks, no explanations.
    **************************FOLLOW ALL THE INSTRUCTION PROPERLY**********************

    OUTPUT FORMAT:

    [
 {{
  "slide_number": 1,
  "slide_title": "Executive Summary",
  "slide_purpose": "Provide board-level overview",

  "content": [
    "- AI became the primary acquisition focus",
    "- Acquisitions expanded cybersecurity capabilities",
    "- Reinvention services accelerated growth"
  ],

  "tables": [],

  "images": [
    {{
      "image_id": "img_001",
      "image_path": "/images/ai_growth.png",
      "image_caption": "Growth of AI-focused acquisitions"
    }}
  ]
}}
    ]

    Return ONLY JSON.
    """

    try:
        response = gemini_client.models.generate_content(
                        model="gemini-3.1-flash-lite",
                        contents=prompt
                    )

        raw_text = response.text.strip()

        
        # raw_text = response.text.strip().replace("\n", "")
        print("XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXx")
        print(raw_text)

    except Exception as e:
         raise e

    slide_plan = json.loads(
                    repair_json_output(raw_text)
                )    
    return {
        "slide_plan": slide_plan
    }
    
def layout_node(slide_plane):

    slide_content = slide_plane

    prompt = f"""
You are a senior executive presentation designer specializing in structural data validation and layout selection.

Your task is to analyze the provided JSON array of slides and select the SINGLE MOST APPROPRIATE layout token for each slide based on its literal text length and active asset metrics.

AVAILABLE LAYOUTS:
1. 'cover_slide'          -> Use ONLY for slide_number: 1 (Title & Subtitle).
2. 'section_divider'     -> Use for transitional slides that introduce a new major narrative block.
3. 'bullet_slide'        -> Use when a slide contains ONLY text/bullet points, and has ZERO tables and ZERO images.
4. 'table_slide'         -> Use when a slide contains a table but has ZERO supporting bullets/text strings.
5. 'image_slide'         -> Use when a slide contains an image/chart path but has ZERO supporting bullets/text strings.
6. 'table_with_insights' -> Use when a slide contains BOTH a valid table AND supporting text/bullet strings.
7. 'image_with_insights' -> Use when a slide contains BOTH a valid image path AND supporting text/bullet strings.
8. 'two_column'          -> Use for general dual-column text balance or multi-entity layouts.
9. 'comparison_slide'    -> Use specifically when the text compares multiple companies, metrics, or periods side-by-side.
10. 'dashboard_slide'    -> High data density layout. Use ONLY when a slide contains BOTH a valid table AND a valid image simultaneously.
11. 'conclusion_slide'   -> Use ONLY for the final slide of the presentation to display recommendations or final takeaways.

STRICT LAYOUT SELECTION RULES:

1. FIRST SLIDE CONSTRAINT:
   - If slide_number == 1, you MUST choose 'cover_slide'.

2. FINAL SLIDE CONSTRAINT:
   - If the slide is the last element in the input array, you MUST choose 'conclusion_slide', regardless of its asset composition.

3. STRICT MIDDLE-SLIDE CONDITIONAL MATRIX (Based on item counts, not key existence):
   - If 'tables' has 0 elements AND 'images' has 0 elements -> 'bullet_slide'
   - If 'tables' has 1+ elements AND 'images' has 0 elements AND 'content' has 1+ elements -> 'table_with_insights'
   - If 'images' has 1+ elements AND 'tables' has 0 elements AND 'content' has 1+ elements -> 'image_with_insights'
   - If 'tables' has 1+ elements AND 'images' has 1+ elements -> 'dashboard_slide'
   - If 'tables' has 1+ elements AND 'content' has 0 elements -> 'table_slide'
   - If 'images' has 1+ elements AND 'content' has 0 elements -> 'image_slide'

4. EMPTY ARRAY WARNING:
   - Do NOT select 'dashboard_slide', 'table_with_insights', or 'image_with_insights' if the corresponding 'tables' or 'images' lists are empty ([ ]). Empty arrays mean the asset does not exist. Fall back to 'bullet_slide'.

INPUT:
{slide_content}

OUTPUT FORMAT

[
  {{
    "slide_number": 1,
    "layout": "cover_slide",
    "reason": "First slide of the presentation"
  }}
]

Return ONLY valid JSON.
"""
    
    try:
        response = gemini_client.models.generate_content(
                        model="gemini-3.1-flash-lite",
                        contents=prompt
                    )

        raw_text = response.text.strip()
        
        # raw_text = response.text.strip().replace("\n", "")
        print("XXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXXx")
        print(raw_text)

    except Exception as e:
         raise e
    
    layout_plan = json.loads(
                    repair_json_output(raw_text)
                )
    
    return {
         "layout_plan": layout_plan
    }

def add_layout_to_slide(state:State):
    # import pdb;pdb.set_trace()
    story_result = story_planner(state["summarized_content"])
    slide_plan = story_result["slide_plan"]

    sanitized_plan_for_llm = []
    for slide in slide_plan:
        clean_slide = slide.copy()
        if "tables" in clean_slide and not clean_slide["tables"]:
            del clean_slide["tables"]
        if "images" in clean_slide and not clean_slide["images"]:
            del clean_slide["images"]
        sanitized_plan_for_llm.append(clean_slide)

    # Get layout recommendations using the clean structure
    layout_result = layout_node(sanitized_plan_for_llm)
    layout_plan = layout_result["layout_plan"]

    # if DEV_MODE and os.path.exists("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/layout.json"):
    #     try: 
    #         with open("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/layout.json") as f:
    #             return {"slide_plan": json.load(f)}
    #     except Exception as e:
    #         print("Invalid cache file. Regenerating...")

    layout_map = {
        item["slide_number"]: item["layout"]
        for item in layout_plan
    }

    # =================================================================
    # STEP 2: PYTHON-LEVEL DEFENSIVE OVERRIDES (THE BULLETPROOF FILTER)
    # =================================================================
    for slide in slide_plan:
        assigned_layout = layout_map.get(slide["slide_number"], "bullet_slide")
        
        # Check if real data actually exists
        has_table = "tables" in slide and len(slide["tables"]) > 0
        has_image = "images" in slide and len(slide["images"]) > 0

        # Override rule: If it chose a layout that needs visuals, but has none, force 'bullet_slide'
        if assigned_layout in ["dashboard_slide", "table_with_insights", "image_with_insights", "table_slide", "image_slide"]:
            if not has_table and not has_image:
                assigned_layout = "bullet_slide"
                
        # Adjust dashboard slide if it only has ONE of the assets instead of both
        if assigned_layout == "dashboard_slide":
            if has_table and not has_image:
                assigned_layout = "table_with_insights"
            elif has_image and not has_table:
                assigned_layout = "image_with_insights"
            elif not has_table and not has_image:
                assigned_layout = "bullet_slide"

        slide["layout"] = assigned_layout

    # Step 3: Explicitly guarantee the last slide is always a conclusion layout
    if slide_plan:
        final_slide_num = max(s["slide_number"] for s in slide_plan)
        for s in slide_plan:
            if s["slide_number"] == final_slide_num:
                s["layout"] = "conclusion_slide"

    print("===================================================================")
    print(slide_plan)

    # with open("/home/sejaldhiman/Desktop/my_folder/Ai-agent/Agents/LangGraph/Agent_langGraph/cache/layout.json", "w") as f:
    #     json.dump(
    #         slide_plan,
    #         f,
    #         indent=2,
    #         ensure_ascii=False
    #     )
    return {
        "slide_plan": slide_plan
    }


def render_cover_slide(slide_data):
    # Safe fallback resolution logic
    title = slide_data.get("slide_title", "Untitled Presentation")
    
    # Fall back to content list if key_message is missing
    subtitle = slide_data.get("key_message", "")
    if not subtitle and slide_data.get("content"):
        subtitle = slide_data["content"][0].replace("-", "")
        
    return f"""
    <div class="h-full flex flex-col justify-center items-start p-32 relative overflow-hidden">
        <div class="w-32 h-3 rounded-full mb-12" style="background-color: var(--accent-color);"></div>
        <h1 class="text-8xl font-bold tracking-tight mb-8 leading-tight text-[var(--primary-color)] max-w-5xl">
            {title}
        </h1>
        {f'<p class="text-3xl font-light leading-relaxed text-[var(--secondary-color)] max-w-3xl">{subtitle}</p>' if subtitle else ''}
    </div>
    """

def parse_markdown_table(md_table_input):
    """
    Safely parses Markdown tables into semantic HTML.
    Handles raw strings, nested 'table_data' dictionaries, and variations in boundary pipe formatting.
    """
    # =========================================================================
    # NEW HANDLING: Process dictionary inputs where 'table_data' is a raw matrix (list of lists)
    # =========================================================================
    if isinstance(md_table_input, dict) and isinstance(md_table_input.get("table_data"), list):
        matrix_rows = md_table_input["table_data"]
        if not matrix_rows:
            return ""
            
        table_html = "<table class='w-full border-collapse text-left my-4'>"
        for idx, row in enumerate(matrix_rows):
            if not isinstance(row, list) or not any(str(cell).strip() for cell in row):
                continue
                
            row_bg = "style='background-color: rgba(14, 165, 233, 0.06);'" if idx % 2 == 0 else "style='background-color: transparent;'"
            table_html += f"<tr class='border-b border-slate-800/60 {row_bg} transition-colors'>"
            
            for cell in row:
                clean_cell = str(cell).replace("**", "").replace("`", "").strip()
                if idx == 0:
                    table_html += f"<th class='p-4 text-xl font-semibold uppercase tracking-wider text-[var(--accent-color)] border-b-2 border-slate-700'>{clean_cell}</th>"
                else:
                    table_html += f"<td class='p-4 text-lg font-light text-[var(--primary-color)] opacity-90'>{clean_cell}</td>"
            table_html += "</tr>"
            
        table_html += "</table>"
        return table_html

    # 1. DEFENSIVE CHECK: Handle dict inputs containing the 'data' key (Preserved)
    if isinstance(md_table_input, dict) and "data" in md_table_input:
        matrix_rows = md_table_input["data"]
        if not isinstance(matrix_rows, list) or not matrix_rows:
            return ""
            
        table_html = "<table class='w-full border-collapse text-left my-4'>"
        for idx, row in enumerate(matrix_rows):
            if not isinstance(row, list) or not any(str(cell).strip() for cell in row):
                continue
                
            row_bg = "bg-slate-900/20" if idx % 2 == 0 else "bg-transparent"
            table_html += f"<tr class='border-b border-slate-800/60 {row_bg} transition-colors'>"
            
            for cell in row:
                clean_cell = str(cell).replace("**", "").replace("`", "").strip()
                if idx == 0:
                    table_html += f"<th class='p-4 text-xl font-semibold uppercase tracking-wider text-[var(--accent-color)] border-b-2 border-slate-700'>{clean_cell}</th>"
                else:
                    table_html += f"<td class='p-4 text-lg font-light text-[var(--primary-color)] opacity-90'>{clean_cell}</td>"
            table_html += "</tr>"
            
        table_html += "</table>"
        return table_html
    
    # (Rest of your unchanged code below for fallback string checks)
    if isinstance(md_table_input, dict):
        md_table_string = md_table_input.get("table_data", "")
    elif isinstance(md_table_input, str):
        md_table_string = md_table_input
    else:
        return ""

    if not md_table_string or not isinstance(md_table_string, str):
        return ""
        
    lines = [line.strip() for line in md_table_string.split("\n") if line.strip() and "---" not in line]
    if not lines:
        return ""
    
    table_html = "<table class='w-full border-collapse text-left my-4'>"
    
    for idx, line in enumerate(lines):
        raw_cells = line.split("|")
        cells = [cell.strip() for cell in raw_cells if cell.strip() or raw_cells.index(cell) not in (0, len(raw_cells)-1)]
        
        if not any(cells):
            continue
            
        row_bg = "bg-slate-900/20" if idx % 2 == 0 else "bg-transparent"
        table_html += f"<tr class='border-b border-slate-800/60 {row_bg} transition-colors'>"
        
        for cell in cells:
            if idx == 0:
                table_html += f"<th class='p-4 text-xl font-semibold uppercase tracking-wider text-[var(--accent-color)] border-b-2 border-slate-700'>{cell}</th>"
            else:
                clean_cell = cell.replace("**", "").replace("`", "").strip()
                table_html += f"<td class='p-4 text-lg font-light text-[var(--primary-color)] opacity-90'>{clean_cell}</td>"
                
        table_html += "</tr>"
        
    table_html += "</table>"
    return table_html

def render_table_insight_slide(slide_data):
    title = slide_data.get("slide_title", "Data Summary")
    raw_table = slide_data["tables"][0] if slide_data.get("tables") else ""
    compiled_table_html = parse_markdown_table(raw_table) if raw_table else "<p class='text-xl text-[var(--secondary-color)]'>No table data provided.</p>"
    
    # Resolve descriptive insight fields safely
    insight_text = slide_data.get("key_message", "")
    if not insight_text and slide_data.get("content"):
        raw_content = slide_data.get("content")
        
        if isinstance(raw_content, list):
            # Clean up leading dashes or markdown dots from each point and join with spaces
            cleaned_points = [item.lstrip("- •").strip() for item in raw_content if item.strip()]
            insight_text = " ".join(cleaned_points)
        elif isinstance(raw_content, str):
            # If it's a raw string block, clean up stray formatting artifacts
            insight_text = raw_content.replace("- ", "").replace("\n", " ").strip()
        
    return f"""
    <div class="h-full grid grid-cols-[60%_40%] gap-16 p-24 items-center">
        <div class="w-full overflow-hidden rounded-2xl border border-slate-800/80 bg-slate-900/30 p-8">
            {compiled_table_html}
        </div>
        <div class="flex flex-col justify-center border-l border-slate-800/80 pl-12">
            <span class="text-lg font-mono uppercase tracking-widest text-[var(--secondary-color)] mb-4">Analytics Matrix</span>
            <h2 class="text-4xl font-bold tracking-tight mb-6 text-[var(--primary-color)]">{title}</h2>
            <p class="text-xl leading-relaxed text-[var(--secondary-color)] font-light max-h-[600px] overflow-hidden">
                {insight_text}
            </p>
        </div>
    </div>
    """

def render_grid_bullet_slide(slide_data):
    title = slide_data.get("slide_title", "Information Analysis")
    items = slide_data.get("content", [])
    
    # Safely handle descriptive messages
    description = slide_data.get("key_message", "")
    icon_class = get_widget_icon(title)
    
    cards_html = ""
    processed_items = []
    temp_header = ""
    
    for item in items:
        cleaned = item.strip()
        if not cleaned:
            continue
            
        # If the item is just a short header ending in a colon, hold onto it
        if cleaned.endswith(":") and len(cleaned) < 25:
            temp_header = f"<strong class='block text-[var(--accent-color)] uppercase tracking-wider text-sm mb-2'>{cleaned}</strong>"
            continue
            
        # If we have a pending header, attach it to the front of this content element
        if temp_header:
            cleaned = f"{temp_header}{cleaned}"
            temp_header = ""
            
        processed_items.append(cleaned)
        
    # If a header was left over at the very end, append it as a standalone point
    if temp_header:
        processed_items.append(temp_header)

    for item in processed_items:
        clean_text = re.sub(r'^([\d+[\.\)]|[\-\•])\s*', '', item)
        cards_html += f"""
        <div class="bg-slate-900/40 border border-slate-800/60 p-8 rounded-2xl flex items-start gap-6 backdrop-blur-md transition-all hover:border-[var(--accent-color)]/40 h-full">
            <div class="flex-shrink-0 w-14 h-14 rounded-xl flex items-center justify-center bg-[var(--accent-color)]/10 border border-[var(--accent-color)]/20">
                 <i class="{icon_class} text-2xl text-[var(--accent-color)]"></i>
            </div>
            <div class="flex-1 flex flex-col justify-start">
                <p class="text-xl font-light leading-relaxed text-[var(--primary-color)] opacity-95">{clean_text}</p>
            </div>
        </div>
        """
        
    # Dynamically scale grid layout columns based on content payload volume
    # grid_cols = "grid-cols-2" if len(items) > 2 else "grid-cols-1"
    grid_cols = "grid-cols-2" if len(items) <= 3 else "grid-cols-3"
    
    return f"""
    <div class="h-full flex flex-col justify-between p-24">
        <div>
            <h2 class="text-7xl font-bold tracking-tight mb-4 text-[var(--primary-color)]">{title}</h2>
            {f'<p class="text-xl font-mono text-[var(--secondary-color)] max-w-4xl font-light">{description}</p>' if description else ''}
        </div>
        <div class="grid {grid_cols} gap-6 flex-1 my-6 overflow-hidden">
            {cards_html}
        </div>
    </div>
    """

def render_conclusion_slide(slide_data):
    return f"""
    <div class="h-full w-full flex flex-col items-center justify-center text-center p-32" 
         style="background-color: var(--inverted-bg); color: #0F172A;">
         
        <span class="text-xl font-mono uppercase tracking-widest text-slate-500 mb-6">Takeaway Summary</span>
        
        <h1 class="text-7xl font-bold tracking-tight mb-10 max-w-5xl leading-tight text-slate-900">
            {slide_data["slide_title"]}
        </h1>
    </div>
    """

def render_image_block(images):
    if not images:
        return ""

    try:
        raw_path = images[0]["image_path"]

        # Convert relative path to absolute path
        absolute_path = os.path.abspath(raw_path)

        print(f"[IMAGE] Raw Path      : {raw_path}")
        print(f"[IMAGE] Absolute Path : {absolute_path}")
        print(f"[IMAGE] Exists        : {os.path.exists(absolute_path)}")

        if not os.path.exists(absolute_path):
            return f"""
            <div class="w-full h-full flex flex-col items-center justify-center">
                <p class="text-red-500 text-xl">
                    Image not found:
                </p>
                <p class="text-slate-400 text-sm mt-2">
                    {absolute_path}
                </p>
            </div>
            """

        # Read image and convert to Base64
        with open(absolute_path, "rb") as img_file:
            encoded = base64.b64encode(
                img_file.read()
            ).decode("utf-8")

        ext = os.path.splitext(absolute_path)[1].lower()

        if ext == ".jpg" or ext == ".jpeg":
            mime = "image/jpeg"
        elif ext == ".webp":
            mime = "image/webp"
        else:
            mime = "image/png"

        browser_safe_url = (
            f"data:{mime};base64,{encoded}"
        )

        caption = images[0].get(
            "image_caption",
            ""
        )

        return f"""
        <div class="w-full h-full flex flex-col items-center justify-center">
            <img
                src="{browser_safe_url}"
                style="
                    max-width:100%;
                    max-height:85%;
                    object-fit:contain;
                "
            />
            <p class="mt-4 text-sm font-mono text-slate-400">
                {caption}
            </p>
        </div>
        """

    except Exception as e:
        print(f"[IMAGE ERROR] {e}")

        return f"""
        <div class="w-full h-full flex items-center justify-center">
            <p class="text-red-500">
                Failed to load image: {str(e)}
            </p>
        </div>
        """

def render_dashboard_slide(slide_data):
    title = slide_data.get("slide_title", "Executive Dashboard")
    bullets = slide_data.get("content", [])
    tables = slide_data.get("tables", [])
    images = slide_data.get("images", [])
    
    # Compile text column
    text_html = ""
    for bullet in bullets:
        text_html += f"""
        <div class="flex items-start gap-4 mb-4 text-xl font-light text-[var(--primary-color)] opacity-90">
            <span class="mt-2.5 w-2 h-2 flex-shrink-0" style="background-color: var(--accent-color);"></span>
            <p>{bullet}</p>
        </div>
        """
    
    # Determine split columns based on data availability
    has_table = len(tables) > 0
    has_image = len(images) > 0
    
    visuals_html = ""
    if has_table and has_image:
        # Split row for side-by-side asset rendering
        visuals_html = f"""
        <div class="grid grid-rows-2 gap-6 h-full">
            <div class="overflow-hidden rounded-xl border border-slate-800 bg-slate-900/40 p-4 text-sm">
                {parse_markdown_table(tables[0])}
            </div>
            <div class="h-full">
                {render_image_block(images)}
            </div>
        </div>
        """
    elif has_table:
        visuals_html = f"""
        <div class="w-full overflow-hidden rounded-xl border border-slate-800 bg-slate-900/40 p-6">
            {parse_markdown_table(tables[0])}
        </div>
        """
    elif has_image:
        visuals_html = render_image_block(images)
    else:
        visuals_html = f"<div class='text-[var(--secondary-color)] font-mono'>No asset payload mapped.</div>"

    return f"""
    <div class="h-full flex flex-col justify-between p-20">
        <div>
            <span class="text-xs font-mono uppercase tracking-widest text-[var(--accent-color)]">Quantitative Metric Performance</span>
            <h1 class="text-5xl font-bold tracking-tight text-[var(--primary-color)] mt-1">{title}</h1>
        </div>
        
        <div class="grid grid-cols-[40%_60%] gap-12 flex-1 mt-8 max-h-[720px] overflow-hidden">
            <div class="flex flex-col justify-start pr-4 border-r border-slate-800/50">
                <h3 class="text-sm font-mono uppercase tracking-wider text-[var(--secondary-color)] mb-6">Strategic Insights</h3>
                <div class="overflow-y-auto max-h-[600px]">
                    {text_html}
                </div>
            </div>
            <div class="h-full overflow-hidden">
                {visuals_html}
            </div>
        </div>
    </div>
    """


def get_widget_icon(slide_title: str) -> str:
    """
    Returns a semantic FontAwesome icon class based on keywords inside the slide title.
    """
    title_lower = slide_title.lower()
    
    # Mapping keywords to clean, corporate geometric icons
    icon_map = {
        "summary": "fa-solid fa-chart-pie",
        "executive": "fa-solid fa-user-tie",
        "market": "fa-solid fa-globe",
        "investment": "fa-solid fa-coins",
        "capital": "fa-solid fa-money-bill-trend-up",
        "financial": "fa-solid fa-scale-balanced",
        "growth": "fa-solid fa-arrow-up-trend-line",
        "risk": "fa-solid fa-triangle-exclamation",
        "tech": "fa-solid fa-microchip",
        "ai": "fa-solid fa-brain",
        "analysis": "fa-solid fa-magnifying-glass-chart",
        "strategy": "fa-solid fa-compass",
        "recommendation": "fa-solid fa-lightbulb",
        "geographic": "fa-solid fa-map-location-dot",
    }
    
    for key, icon in icon_map.items():
        if key in title_lower:
            return icon
            
    # Default clean widget fallback icon if no keyword matches
    return "fa-solid fa-cubes"



def slide_renderer_node(state: State):
    # import pdb;pdb.set_trace()
    slides = state["slide_plan"]
    design_system = state["design_config"]

    print("=====================================================================================")
    print(slides)
    
    # 1. Standard CSS Injection Setup matching your Technical_Modern tokens
    css_variables = f"""
    <style>
        :root {{
            --bg-canvas: {design_system["bg_canvas"]};
            --primary-color: {design_system["primary"]};
            --secondary-color: {design_system["secondary"]};
            --accent-color: {design_system["accent"]};
            --inverted-bg: {design_system["inverted_section_bg"]};
        }}
        .slide-container {{
            width: 1920px;
            height: 1080px;
            background-color: var(--bg-canvas);
            box-sizing: border-box;
            overflow: hidden;
            position: relative;
        }}
        .style-tech_grid {{
            background-image: 
                linear-gradient(to right, rgba(100, 116, 139, 0.04) 1px, transparent 2px),
                linear-gradient(to bottom, rgba(100, 116, 139, 0.04) 1px, transparent 2px);
            background-size: 50px 50px;
        }}
    </style>
    """
    
    output_images = []
    
    # 2. Continuous Playwright Thread Process
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page()
        page.set_viewport_size({"width": 1920, "height": 1080})

        print(f"[DEBUG] Total incoming slides to process: {len(slides)}")
        
        for idx, slide in enumerate(slides):
            print(f"[DEBUG] Processing loop pass for slide index: {idx + 1}")
            layout_type = slide.get("layout", "bullet_slide")

            if layout_type == "cover_slide":
                body_content = render_cover_slide(slide)
            elif layout_type == "dashboard_slide":
                body_content = render_dashboard_slide(slide)
            elif layout_type in ["table_with_insights", "table_slide"]:
                body_content = render_table_insight_slide(slide)
            elif layout_type in ["image_slide", "image_with_insights"]:
                # Reuse your component logic safely
                slide["tables"] = [] # Force view context to prioritize image blocks
                body_content = render_dashboard_slide(slide) 
            else:
                body_content = render_grid_bullet_slide(slide)
                
            # Compile integrated layout wrapper frame
            # The footer brand elements are completely hidden on full-bleed inverted slides
            show_footer = "flex" if layout_type != "conclusion_slide" else "hidden"
            
            full_html = f"""
            <!DOCTYPE html>
            <html>
            <head>
                <script src="https://cdn.tailwindcss.com"></script>
                <link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/font-awesome/6.4.0/css/all.min.css">
                {css_variables}
            </head>
            <body class="m-0 p-0">
                <div class="slide-container style-{design_system["background_style"]}">
                    {body_content}
                    
                    <div class="absolute bottom-12 left-24 right-24 {show_footer} justify-between items-center text-sm font-mono text-[var(--secondary-color)]">
                        <span>PAGE {str(idx + 1).zfill(2)}</span>
                    </div>
                </div>
            </body>
            </html>
            """
            
            page.set_content(full_html)

            page.wait_for_load_state("networkidle")

            page.wait_for_timeout(1000)# Brief structural rendering catch
            
            filepath = f"slide_deck_page_{idx + 1}.png"
            page.screenshot(path=filepath)
            output_images.append(filepath)
            print(f"[DEBUG] Total images appended to array: {len(output_images)}") 
        browser.close()
        
    return {"rendered_slide_images": output_images}

def export_to_pptx_node(state: State):
    image_paths = state["rendered_slide_images"]
    
    prs = Presentation()
    # Set slide dimensions to Widescreen (16:9)
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for img_path in image_paths:
        # Add a blank slide layout
        blank_slide_layout = prs.slide_layouts[6] 
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Add the PNG image to fill the entire slide
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    pptx_filename = "Final_Presentation.pptx"
    prs.save(pptx_filename)
    
    print(f"Success: {pptx_filename} generated with {len(image_paths)} slides.")
    return {"pptx_path": pptx_filename}